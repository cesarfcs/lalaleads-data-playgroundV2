import io
from datetime import datetime
import pandas as pd
import numpy as np
import streamlit as st
import plotly.express as px
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import cm
from reportlab.pdfgen import canvas
from PIL import Image

# -----------------------
# Branding Lalaleads
# -----------------------
BRAND_BG = "#000000"
BRAND_PRIMARY = "#852e31"
BRAND_TEXT = "#fcfcfc"
BRAND_MUTED = "#a7a7a7"
GRID_COLOR = "#222222"

st.set_page_config(page_title="Lalaleads – Data Playground HubSpot", layout="wide")

st.markdown(f"""
<style>
html, body, [class^="css"] {{ background-color: {BRAND_BG} !important; color: {BRAND_TEXT} !important; }}
section[data-testid="stSidebar"] {{ background: {BRAND_BG}; border-right: 1px solid {GRID_COLOR}; }}
h1,h2,h3,h4 {{ color: {BRAND_TEXT} !important; }}
div[data-testid="stMetric"] {{ background: rgba(255,255,255,0.02); border: 1px solid {GRID_COLOR}; border-radius: 14px; padding: 14px; }}
button[kind="primary"] {{ background: {BRAND_PRIMARY} !important; color: {BRAND_TEXT} !important; border-radius: 10px; }}
button {{ border-radius: 10px; }}
</style>
""", unsafe_allow_html=True)

# Header & logo
try:
    st.image("assets/logo.png", width=220)
except Exception:
    pass
st.title("🧩 Lalaleads – Data Playground (HubSpot)")
st.caption("Appels, emails, RDV, secteurs et intitulés de poste – filtres multi-critères, KPIs auto, exports PDF/PPTX.")

# Plotly theme
PLOTLY_TEMPLATE = dict(
    layout=dict(
        font=dict(color=BRAND_TEXT),
        paper_bgcolor=BRAND_BG,
        plot_bgcolor=BRAND_BG,
        legend=dict(bgcolor="rgba(0,0,0,0)", bordercolor=GRID_COLOR),
        xaxis=dict(gridcolor=GRID_COLOR, zerolinecolor=GRID_COLOR),
        yaxis=dict(gridcolor=GRID_COLOR, zerolinecolor=GRID_COLOR),
    )
)
COLOR_SEQUENCE = [BRAND_PRIMARY, BRAND_TEXT, BRAND_MUTED]
POSSIBLE_DELIMS = [",", ";", "\t", "|"]

@st.cache_data(show_spinner=False)
def read_csv_smart(file, encoding="utf-8", delimiter=None):
    file.seek(0)
    try:
        return pd.read_csv(file, encoding=encoding, sep=delimiter)
    except Exception:
        pass
    for d in POSSIBLE_DELIMS:
        file.seek(0)
        try:
            return pd.read_csv(file, encoding=encoding, sep=d)
        except Exception:
            continue
    file.seek(0)
    try:
        return pd.read_excel(file)
    except Exception as e:
        raise e

# -----------------------
# Mapping hints
# -----------------------
HINTS = {
    "date": ["date", "created", "last", "activity", "appel", "email"],
    "client": ["client", "company", "entreprise", "account", "nom de l'entreprise"],
    "campaign": ["campagne", "campaign"],
    "bdr": ["commercial", "owner", "propriétaire", "bdr", "sales"],
    "lifecycle": ["cycle", "phase", "lifecycle", "rdv"],
    "aircall": ["aircall", "call", "appel", "tags"],
    "lemlist": ["lemlist", "email", "status", "statut"],
    "contacts_count": ["prises de contact", "touches", "attempts", "essais", "nombre de prises"],
    "job_title": ["intitulé", "poste", "title", "fonction"],
    "industry": ["secteur", "industrie", "industry"],
}

def guess_column(cols, hints):
    cols_lower = {c.lower(): c for c in cols}
    for h in hints:
        for c in cols_lower:
            if h in c:
                return cols_lower[c]
    return None

# -----------------------
# Upload & mapping
# -----------------------
st.sidebar.header("⬆️ Import HubSpot CSV/XLSX")
files = st.sidebar.file_uploader(
    "Glissez-déposez 1+ fichiers (exports HubSpot)",
    type=["csv", "xlsx"],
    accept_multiple_files=True,
)
encoding = st.sidebar.selectbox("Encodage", ["utf-8", "latin-1", "cp1252"], index=0)

logo_file = st.sidebar.file_uploader("Logo (pour exports)", type=["png","jpg","jpeg"])
logo_img = None
if logo_file is not None:
    logo_img = Image.open(logo_file).convert("RGBA")
else:
    try:
        logo_img = Image.open("assets/logo.png").convert("RGBA")
    except Exception:
        logo_img = None

if not files:
    st.info("Chargez vos exports HubSpot pour commencer (ou utilisez un CSV d'exemple).")
    st.stop()

raw_dfs = []
for f in files:
    df = read_csv_smart(f, encoding=encoding)
    df["__source_file"] = f.name
    raw_dfs.append(df)

raw = pd.concat(raw_dfs, ignore_index=True)
raw_cols = list(raw.columns)

st.sidebar.markdown("---")
st.sidebar.subheader("🧩 Mapping des colonnes (HubSpot)")

def autoselect(label, key):
    guess = guess_column(raw_cols, HINTS[key])
    idx = (raw_cols.index(guess) + 1) if (guess is not None and guess in raw_cols) else 0
    return st.sidebar.selectbox(label, [None] + raw_cols, index=idx)

col_date     = autoselect("Colonne date d'activité (appels/emails)", "date")
col_campaign = autoselect("Campagne", "campaign")
col_bdr      = autoselect("Commercial (BDR)", "bdr")
col_lifecycle= autoselect("Phase du cycle de vie", "lifecycle")
col_aircall  = autoselect("Statut / Tags Appel (Aircall)", "aircall")
col_lemlist  = autoselect("Statut Email (lemlist)", "lemlist")
col_contacts = autoselect("Nombre de prises de contact (tentatives)", "contacts_count")
col_client   = autoselect("Client / Entreprise", "client")
col_job      = autoselect("Intitulé de poste", "job_title")
col_industry = autoselect("Secteur d'activité", "industry")

work = raw.copy()

# Dates
if col_date:
    with st.spinner("Normalisation des dates…"):
        work[col_date] = pd.to_datetime(work[col_date], errors='coerce')

# Numériques
def to_num(s):
    if s is None: return np.nan
    try: return pd.to_numeric(str(s).replace(" ", "").replace(",", "."), errors='coerce')
    except Exception: return np.nan
if col_contacts and col_contacts in work.columns:
    work[col_contacts] = work[col_contacts].apply(to_num)

# -----------------------
# Filtres
# -----------------------
st.sidebar.subheader("🔎 Filtres")
if col_date and work[col_date].dropna().size > 0:
    min_d = pd.to_datetime(work[col_date]).min().date()
    max_d = pd.to_datetime(work[col_date]).max().date()
    d_from, d_to = st.sidebar.date_input("Période", value=(min_d, max_d), min_value=min_d, max_value=max_d)
    mask = (work[col_date] >= pd.to_datetime(d_from)) & (work[col_date] <= pd.to_datetime(d_to) + pd.Timedelta(days=1) - pd.Timedelta(seconds=1))
    work = work.loc[mask]

def multisel(col, label, default_n=10):
    if col and col in work.columns:
        opts = sorted([x for x in work[col].dropna().astype(str).unique()])
        chosen = st.sidebar.multiselect(label, opts, default=opts[:min(len(opts), default_n)])
        if chosen: return work[work[col].astype(str).isin(chosen)]
    return work

work = multisel(col_campaign, "Campagnes", 10)
work = multisel(col_bdr, "BDR", 10)
work = multisel(col_job, "Intitulés de poste", 10)
work = multisel(col_industry, "Secteurs d'activité", 10)

st.sidebar.markdown("---")
report_title = st.sidebar.text_input("Titre du rapport", value=f"Rapport – {datetime.now().strftime('%Y-%m-%d')}")
company_name = st.sidebar.text_input("Nom du client / mission", value="Lalaleads – Mission")
notes_field = st.sidebar.text_area("Remontées terrain (optionnel)")
next_steps_field = st.sidebar.text_area("Next steps (optionnel)")
btn_pptx = st.sidebar.button("📤 Exporter en PowerPoint (.pptx)")
btn_pdf  = st.sidebar.button("📄 Exporter en PDF (.pdf)")

# -----------------------
# KPIs
# -----------------------
def contains_any(val, keywords):
    if pd.isna(val): return False
    s = str(val).lower()
    return any(k in s for k in keywords)

CALL_CONNECTED = ["connected", "décro", "answered", "meeting", "rappel", "callback", "qualifié", "qualifie"]
CALL_PITCHED   = ["pitch", "demo", "présentation", "presentation"]
CALL_REFUSAL   = ["refus", "pas intéress", "not interested", "désint"]
EMAIL_OPENED   = ["opened", "open"]
EMAIL_REPLIED  = ["replied", "répon", "reply", "interested"]
RDV_KEYWORDS   = ["rdv", "meeting", "rendez-vous"]

total_contacts = int(work.shape[0])
appels_passes  = int(work[col_contacts].sum()) if (col_contacts and col_contacts in work.columns) else total_contacts

def count_if(series, keywords):
    if series is None or len(series) == 0: return 0
    return int(series.apply(lambda x: contains_any(x, keywords)).sum())

aircall_series  = work[col_aircall]   if (col_aircall   and col_aircall   in work.columns) else pd.Series([], dtype=object)
lemlist_series  = work[col_lemlist]   if (col_lemlist   and col_lemlist   in work.columns) else pd.Series([], dtype=object)
lifecycle_series= work[col_lifecycle] if (col_lifecycle and col_lifecycle in work.columns) else pd.Series([], dtype=object)

connectes = count_if(aircall_series, CALL_CONNECTED)
pitch     = count_if(aircall_series, CALL_PITCHED)
refus     = count_if(aircall_series, CALL_REFUSAL)
rdv       = count_if(lifecycle_series, RDV_KEYWORDS) + count_if(aircall_series, RDV_KEYWORDS)
opened    = count_if(lemlist_series, EMAIL_OPENED)
replied   = count_if(lemlist_series, EMAIL_REPLIED)
sent_emails = total_contacts  # fallback si pas de colonne dédiée

def rate(a, b): return (a / b * 100) if b and b > 0 else 0.0

r1 = st.columns(4)
r1[0].metric("Appels passés", f"{appels_passes}")
r1[1].metric("Taux de connexion", f"{rate(connectes, appels_passes):.1f}%")
r1[2].metric("Taux de pitch", f"{rate(pitch, max(connectes,1)):.1f}%")
r1[3].metric("Taux de conversion RDV", f"{rate(rdv, max(pitch or connectes,1)):.1f}%")

r2 = st.columns(3)
r2[0].metric("Taux de désintérêt", f"{rate(refus, max(connectes,1)):.1f}%")
r2[1].metric("Taux d'ouverture email", f"{rate(opened, sent_emails):.1f}%")
r2[2].metric("Taux de réponse email", f"{rate(replied, max(opened,1)):.1f}%")

st.markdown("---")

# -----------------------
# Graphiques
# -----------------------
charts = {}

# RDV par campagne
if col_campaign and col_campaign in work.columns:
    tmp = work.copy()
    tmp['__is_rdv'] = False
    if col_lifecycle and col_lifecycle in tmp.columns:
        tmp['__is_rdv'] = tmp['__is_rdv'] | tmp[col_lifecycle].apply(lambda x: contains_any(x, RDV_KEYWORDS))
    if col_aircall and col_aircall in tmp.columns:
        tmp['__is_rdv'] = tmp['__is_rdv'] | tmp[col_aircall].apply(lambda x: contains_any(x, RDV_KEYWORDS))
    group_col = col_client if (col_client and col_client in tmp.columns) else '__source_file'
    grp = tmp.groupby(col_campaign).agg(contacts=(group_col,'count'), rdv=('__is_rdv','sum')).reset_index()
    grp['taux_rdv_%'] = grp['rdv'] / grp['contacts'] * 100
    fig_conv_camp = px.bar(grp.sort_values('taux_rdv_%', ascending=False), x=col_campaign, y='taux_rdv_%',
                           title='Taux de conversion RDV par campagne', template=PLOTLY_TEMPLATE,
                           color_discrete_sequence=[BRAND_PRIMARY])
    st.plotly_chart(fig_conv_camp, use_container_width=True)
    charts["Conversion RDV par campagne"] = fig_conv_camp

# Connexion par secteur
if col_industry and col_industry in work.columns:
    grp = work.groupby(col_industry).size().reset_index(name='contacts')
    if col_aircall and col_aircall in work.columns:
        grp_conn = work.groupby(col_industry)[col_aircall].apply(lambda s: sum(s.apply(lambda x: contains_any(x, CALL_CONNECTED)))).reset_index(name='connectes')
        grp = grp.merge(grp_conn, on=col_industry, how='left')
        grp['taux_conn_%'] = grp['connectes'] / grp['contacts'] * 100
        fig_sector = px.bar(grp.sort_values('taux_conn_%', ascending=False), x=col_industry, y='taux_conn_%',
                            title="Taux de connexion par secteur d'activité", template=PLOTLY_TEMPLATE,
                            color_discrete_sequence=[BRAND_PRIMARY])
        st.plotly_chart(fig_sector, use_container_width=True)
        charts["Connexion par secteur"] = fig_sector

# Réponse email par intitulé de poste
if col_job and col_job in work.columns and col_lemlist and col_lemlist in work.columns:
    grp = work.groupby(col_job).size().reset_index(name='contacts')
    grp_rep = work.groupby(col_job)[col_lemlist].apply(lambda s: sum(s.apply(lambda x: contains_any(x, EMAIL_REPLIED)))).reset_index(name='replies')
    grp = grp.merge(grp_rep, on=col_job, how='left')
    grp['taux_rep_%'] = grp['replies'] / grp['contacts'] * 100
    fig_role = px.bar(grp.sort_values('taux_rep_%', ascending=False).head(20), x=col_job, y='taux_rep_%',
                      title="Taux de réponse email par intitulé de poste", template=PLOTLY_TEMPLATE,
                      color_discrete_sequence=[BRAND_PRIMARY])
    st.plotly_chart(fig_role, use_container_width=True)
    charts["Réponse par poste"] = fig_role

# Donut emails
if col_lemlist and col_lemlist in work.columns:
    donut = pd.DataFrame({'statut': ['Ouverts','Répondus','Non ouverts'], 'val': [opened, replied, max(sent_emails - opened, 0)]})
    fig_donut = px.pie(donut, values='val', names='statut', hole=0.5, title='Répartition emails',
                       template=PLOTLY_TEMPLATE, color_discrete_sequence=COLOR_SEQUENCE)
    st.plotly_chart(fig_donut, use_container_width=True)
    charts["Répartition emails"] = fig_donut

# RDV par BDR (optionnel)
if col_bdr and col_bdr in work.columns:
    tmp = work.copy()
    tmp['__is_rdv'] = False
    if col_lifecycle and col_lifecycle in tmp.columns:
        tmp['__is_rdv'] = tmp['__is_rdv'] | tmp[col_lifecycle].apply(lambda x: contains_any(x, RDV_KEYWORDS))
    if col_aircall and col_aircall in tmp.columns:
        tmp['__is_rdv'] = tmp['__is_rdv'] | tmp[col_aircall].apply(lambda x: contains_any(x, RDV_KEYWORDS))
    grp = tmp.groupby(col_bdr)['__is_rdv'].sum().reset_index(name='rdv')
    fig_bdr = px.bar(grp.sort_values('rdv', ascending=False), x=col_bdr, y='rdv',
                     title='RDV par BDR', template=PLOTLY_TEMPLATE, color_discrete_sequence=[BRAND_PRIMARY])
    st.plotly_chart(fig_bdr, use_container_width=True)
    charts["RDV par BDR"] = fig_bdr

st.markdown("### Données filtrées")
st.dataframe(work, use_container_width=True)

@st.cache_data(show_spinner=False)
def fig_to_png_bytes(fig):
    try:
        return fig.to_image(format="png", scale=2)  # nécessite kaleido
    except Exception as e:
        st.error("Export d'images indisponible (kaleido manquant). Vérifie requirements.txt.")
        raise e

# -----------------------
# Exports PDF / PPTX
# -----------------------
def export_pptx(title, company, charts_dict, kpis, notes, next_steps, logo=None):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = company
    if logo is not None:
        bio_logo = io.BytesIO(); logo.save(bio_logo, format="PNG"); bio_logo.seek(0)
        slide.shapes.add_picture(bio_logo, Inches(9.0), Inches(0.2), height=Inches(0.6))
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "KPIs"
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4.5)).text_frame
    tx.word_wrap = True
    for label, value in kpis:
        run = tx.paragraphs[0].add_run(); run.text = f"• {label}: {value}\n"; run.font.size = Pt(18)
    for name, fig in charts_dict.items():
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = name
        png = fig_to_png_bytes(fig); image_stream = io.BytesIO(png)
        slide.shapes.add_picture(image_stream, Inches(0.5), Inches(1.5), width=Inches(9))
    if notes or next_steps:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Notes & Next steps"
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4.5)).text_frame
        if notes: tx.text = f"Remontées terrain\n{notes}"
        if next_steps:
            p = tx.add_paragraph(); p.text = f"\nNext steps\n{next_steps}"
    bio = io.BytesIO(); prs.save(bio); bio.seek(0); return bio

def export_pdf(title, company, charts_dict, kpis, notes, next_steps, logo=None):
    buffer = io.BytesIO(); c = canvas.Canvas(buffer, pagesize=A4); width, height = A4
    c.setFillColorRGB(0,0,0); c.rect(0,0,width,height, fill=1, stroke=0)
    c.setFillColorRGB(1,1,1); c.setFont("Helvetica-Bold", 22); c.drawString(2*cm, height-3*cm, title)
    c.setFont("Helvetica", 12); c.drawString(2*cm, height-4*cm, company)
    r,g,b = tuple(int(BRAND_PRIMARY.strip('#')[i:i+2], 16)/255 for i in (0,2,4))
    c.setFillColorRGB(r,g,b); c.rect(2*cm, height-4.2*cm, width-4*cm, 3, fill=1, stroke=0)
    if logo is not None:
        bio_logo = io.BytesIO(); logo.save(bio_logo, format="PNG"); bio_logo.seek(0)
        c.drawImage(bio_logo, width-5*cm, height-3.5*cm, 3.5*cm, 1.2*cm, mask='auto')
    c.showPage()
    c.setFillColorRGB(0,0,0); c.rect(0,0,width,height, fill=1, stroke=0)
    c.setFillColorRGB(1,1,1); c.setFont("Helvetica-Bold", 16); c.drawString(2*cm, height-2.5*cm, "KPIs")
    c.setFont("Helvetica", 12); y = height-4*cm
    for label, value in kpis: c.drawString(2*cm, y, f"• {label}: {value}"); y -= 0.8*cm
    c.showPage()
    for name, fig in charts_dict.items():
        png = fig_to_png_bytes(fig); im = Image.open(io.BytesIO(png))
        c.setFillColorRGB(1,1,1); c.rect(0,0,width,height, fill=1, stroke=0)
        c.setFillColorRGB(0,0,0); c.setFont("Helvetica-Bold", 14); c.drawString(2*cm, height-2.0*cm, name)
        max_w, max_h = width-4*cm, height-6*cm; iw, ih = im.size; scale = min(max_w/iw, max_h/ih); w, h = iw*scale, ih*scale
        c.drawInlineImage(im, (width-w)/2, (height-h)/2 - 0.5*cm, w, h); c.showPage()
    if notes or next_steps:
        c.setFillColorRGB(0,0,0); c.rect(0,0,width,height, fill=1, stroke=0)
        c.setFillColorRGB(1,1,1); c.setFont("Helvetica-Bold", 16); c.drawString(2*cm, height-2.5*cm, "Notes & Next steps")
        c.setFont("Helvetica", 12); y = height-4*cm
        if notes:
            for line in str(notes).split("\n"): c.drawString(2*cm, y, line); y -= 0.6*cm
        if next_steps:
            y -= 0.6*cm; c.setFont("Helvetica-Bold", 12); c.drawString(2*cm, y, "Next steps"); y -= 0.6*cm; c.setFont("Helvetica", 12)
            for line in str(next_steps).split("\n"): c.drawString(2*cm, y, line); y -= 0.6*cm
        c.showPage()
    c.save(); buffer.seek(0); return buffer

kpis = [
    ("Appels passés", f"{appels_passes}"),
    ("Taux de connexion", f"{rate(connectes, appels_passes):.1f}%"),
    ("Taux de pitch", f"{rate(pitch, max(connectes,1)):.1f}%"),
    ("Taux de conversion RDV", f"{rate(rdv, max(pitch or connectes,1)):.1f}%"),
    ("Taux de désintérêt", f"{rate(refus, max(connectes,1)):.1f}%"),
    ("Taux d'ouverture email", f"{rate(opened, sent_emails):.1f}%"),
    ("Taux de réponse email", f"{rate(replied, max(opened,1)):.1f}%"),
]

c1, c2 = st.columns(2)
with c1:
    if btn_pptx:
        bio = export_pptx(report_title, company_name, charts, kpis, notes_field, next_steps_field, logo=logo_img)
        st.download_button("Télécharger le rapport PowerPoint", data=bio, file_name=f"{report_title.replace(' ', '_')}.pptx",
                           mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", use_container_width=True)
with c2:
    if btn_pdf:
        pdf = export_pdf(report_title, company_name, charts, kpis, notes_field, next_steps_field, logo=logo_img)
        st.download_button("Télécharger le rapport PDF", data=pdf, file_name=f"{report_title.replace(' ', '_')}.pdf",
                           mime="application/pdf", use_container_width=True)

with st.expander("❓ Conseils de mapping HubSpot"):
    st.write("""
    - **Date d'activité**: idéalement `Last activity date` ou `Date d'appel / email`.
    - **Statut Appel (Aircall)**: tags type `No answer`, `Meeting`, `Callback`, `Pitch`…
    - **Statut Email (lemlist)**: `Email opened`, `Email replied`, `Interested`…
    - **Cycle de vie**: utilisé pour détecter les `RDV` (si libellé présent).
    - **Intitulé de poste** & **Secteur d'activité**: activent les graphiques correspondants.
    - **Prises de contact**: si absent, l'app estime via le volume de lignes.
    """)
